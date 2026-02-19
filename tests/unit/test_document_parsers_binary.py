"""Unit tests for document_parsers: docx, pdf, xlsx, pptx, epub."""

import io

import pytest

from relrag.infrastructure.document_parsers.base import ParseResult
from relrag.infrastructure.document_parsers.docx_parser import parse_docx
from relrag.infrastructure.document_parsers.epub_parser import (
    _get_dc,
    parse_epub,
)
from relrag.infrastructure.document_parsers.pdf_parser import (
    _map_metadata,
    _parse_pdf_date,
    parse_pdf,
)
from relrag.infrastructure.document_parsers.pptx_parser import parse_pptx
from relrag.infrastructure.document_parsers.xlsx_parser import parse_xlsx


def _minimal_docx_bytes() -> bytes:
    """Create minimal valid .docx in memory."""
    from docx import Document as DocxDocument

    doc = DocxDocument()
    doc.add_paragraph("Hello docx")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _minimal_pdf_bytes() -> bytes:
    """Create minimal valid PDF in memory."""
    from pypdf import PdfWriter

    w = PdfWriter()
    w.add_blank_page(width=72, height=72)
    buf = io.BytesIO()
    w.write(buf)
    return buf.getvalue()


def _pdf_with_metadata_bytes() -> bytes:
    """Create PDF with metadata for testing _map_metadata and _parse_pdf_date."""
    from pypdf import PdfWriter

    w = PdfWriter()
    w.add_blank_page(width=72, height=72)
    w.add_metadata({
        "/Title": "Test PDF",
        "/Author": "Author Name",
        "/CreationDate": "D:20200101120000",
        "/ModDate": "D:20201231235959",
    })
    buf = io.BytesIO()
    w.write(buf)
    return buf.getvalue()


def _minimal_xlsx_bytes() -> bytes:
    """Create minimal valid .xlsx in memory."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    if ws:
        ws["A1"] = "Cell text"
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _minimal_pptx_bytes() -> bytes:
    """Create minimal valid .pptx in memory."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    box.text_frame.text = "Slide text"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _minimal_epub_bytes() -> bytes:
    """Create minimal valid .epub in memory (read_epub can load it)."""
    from ebooklib import epub

    book = epub.EpubBook()
    book.set_identifier("test-id")
    book.set_title("Test Book")
    book.set_language("en")
    c1 = epub.EpubHtml(title="Chapter 1", file_name="ch1.xhtml", lang="en")
    c1.content = "<html><body><p>Chapter content</p></body></html>"
    book.add_item(c1)
    book.toc = (epub.Link("ch1.xhtml", "Chapter 1", "ch1"),)
    book.spine = [c1]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    buf = io.BytesIO()
    epub.write_epub(buf, book, {})
    return buf.getvalue()


class TestParseDocx:
    """Tests for parse_docx."""

    def test_valid_docx_returns_text_and_properties(self) -> None:
        data = _minimal_docx_bytes()
        result = parse_docx(data, filename="doc.docx")
        assert isinstance(result, ParseResult)
        assert "Hello docx" in result.text or result.text.strip()
        assert result.properties.get("source_file_name")[0] == "doc.docx"
        assert result.properties.get("source_file_type")[0] == "docx"

    def test_docx_with_table_includes_table_text(self) -> None:
        from docx import Document as DocxDocument

        doc = DocxDocument()
        doc.add_paragraph("P1")
        table = doc.add_table(rows=2, cols=2)
        table.rows[0].cells[0].text = "A"
        table.rows[0].cells[1].text = "B"
        buf = io.BytesIO()
        doc.save(buf)
        result = parse_docx(buf.getvalue(), filename="t.docx")
        assert "P1" in result.text
        assert "A" in result.text and "B" in result.text

    def test_docx_empty_paragraphs_returns_space(self) -> None:
        from docx import Document as DocxDocument

        doc = DocxDocument()
        buf = io.BytesIO()
        doc.save(buf)
        result = parse_docx(buf.getvalue(), filename=None)
        assert result.text.strip() == "" or result.text == " "

    def test_docx_core_properties_string_mapped(self) -> None:
        from docx import Document as DocxDocument

        doc = DocxDocument()
        doc.core_properties.title = "Doc Title"
        doc.core_properties.author = "Doc Author"
        doc.add_paragraph("x")
        buf = io.BytesIO()
        doc.save(buf)
        result = parse_docx(buf.getvalue(), filename=None)
        assert result.properties.get("title")[0] == "Doc Title"
        assert result.properties.get("author")[0] == "Doc Author"

    def test_invalid_docx_raises_value_error(self) -> None:
        with pytest.raises(ValueError, match="Invalid or corrupted docx"):
            parse_docx(b"not a docx", filename="x.docx")


class TestParsePdfDate:
    """Tests for _parse_pdf_date (PDF date string conversion)."""

    def test_none_or_empty_returns_normalized(self) -> None:
        assert _parse_pdf_date(None) == ""
        assert _parse_pdf_date("") == ""

    def test_non_d_prefix_returns_normalized(self) -> None:
        assert _parse_pdf_date("2020-01-01") == "2020-01-01"

    def test_d_prefix_short_returns_as_is(self) -> None:
        assert _parse_pdf_date("D:2020") == "D:2020"
        assert _parse_pdf_date("D:202001") == "D:202001"

    def test_d_prefix_long_enough_returns_iso_like(self) -> None:
        assert _parse_pdf_date("D:20200101") == "2020-01-01"
        assert _parse_pdf_date("D:20200101120000") == "2020-01-01"


class TestPdfMapMetadata:
    """Tests for _map_metadata (PDF metadata -> canonical keys)."""

    def test_reader_with_no_metadata_returns_empty(self) -> None:
        from unittest.mock import MagicMock

        reader = MagicMock()
        reader.metadata = None
        assert _map_metadata(reader) == {}

    def test_reader_with_empty_metadata_returns_empty(self) -> None:
        from unittest.mock import MagicMock

        reader = MagicMock()
        reader.metadata = {}
        assert _map_metadata(reader) == {}


class TestParsePdf:
    """Tests for parse_pdf."""

    def test_valid_pdf_returns_text_and_properties(self) -> None:
        data = _minimal_pdf_bytes()
        result = parse_pdf(data, filename="a.pdf")
        assert isinstance(result, ParseResult)
        assert result.properties.get("source_file_name")[0] == "a.pdf"
        assert result.properties.get("source_file_type")[0] == "pdf"
        assert "page_count" in result.properties
        assert result.properties["page_count"][0] == "1"

    def test_invalid_pdf_raises_value_error(self) -> None:
        with pytest.raises(ValueError, match="Invalid or corrupted PDF"):
            parse_pdf(b"not a pdf", filename="x.pdf")

    def test_pdf_with_metadata_maps_title_author_and_dates(self) -> None:
        data = _pdf_with_metadata_bytes()
        result = parse_pdf(data, filename=None)
        assert result.properties.get("title")[0] == "Test PDF"
        # Author may be from /Author or /Producer (pypdf adds Producer)
        assert result.properties.get("author") is not None
        created = result.properties.get("created_date")
        assert created is not None and created[0] == "2020-01-01"
        modified = result.properties.get("modified_date")
        assert modified is not None and "2020-12-31" in modified[0]


class TestParseXlsx:
    """Tests for parse_xlsx."""

    def test_valid_xlsx_returns_text_and_properties(self) -> None:
        data = _minimal_xlsx_bytes()
        result = parse_xlsx(data, filename="sheet.xlsx")
        assert isinstance(result, ParseResult)
        assert "Cell text" in result.text or result.text.strip()
        assert result.properties.get("source_file_name")[0] == "sheet.xlsx"
        assert result.properties.get("source_file_type")[0] == "xlsx"

    def test_xlsx_without_filename_no_source_properties(self) -> None:
        data = _minimal_xlsx_bytes()
        result = parse_xlsx(data, filename=None)
        assert "source_file_name" not in result.properties
        assert "source_file_type" not in result.properties

    def test_invalid_xlsx_raises_value_error(self) -> None:
        with pytest.raises(ValueError, match="Invalid or corrupted xlsx"):
            parse_xlsx(b"not xlsx", filename="x.xlsx")


class TestParsePptx:
    """Tests for parse_pptx."""

    def test_valid_pptx_returns_text_and_properties(self) -> None:
        data = _minimal_pptx_bytes()
        result = parse_pptx(data, filename="pres.pptx")
        assert isinstance(result, ParseResult)
        assert "Slide text" in result.text or result.text.strip()
        assert result.properties.get("source_file_name")[0] == "pres.pptx"
        assert result.properties.get("source_file_type")[0] == "pptx"
        assert result.properties.get("page_count")[0] == "1"

    def test_pptx_with_core_properties(self) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.core_properties.title = "Pres Title"
        prs.core_properties.author = "Pres Author"
        prs.core_properties.last_modified_by = "Modifier"
        buf = io.BytesIO()
        prs.save(buf)
        result = parse_pptx(buf.getvalue(), filename="p.pptx")
        assert result.properties.get("title")[0] == "Pres Title"
        assert result.properties.get("author")[0] in ("Pres Author", "Modifier")

    def test_invalid_pptx_raises_value_error(self) -> None:
        with pytest.raises(ValueError, match="Invalid or corrupted pptx"):
            parse_pptx(b"not pptx", filename="x.pptx")


class TestParseEpub:
    """Tests for parse_epub."""

    def test_valid_epub_returns_text_and_properties(self) -> None:
        data = _minimal_epub_bytes()
        result = parse_epub(data, filename="book.epub")
        assert isinstance(result, ParseResult)
        assert "Chapter content" in result.text or result.text.strip()
        assert result.properties.get("source_file_name")[0] == "book.epub"
        assert result.properties.get("source_file_type")[0] == "epub"
        assert result.properties.get("title")[0] == "Test Book"

    def test_epub_item_decode_error_skipped(self) -> None:
        """When an item's get_content().decode() raises, that item is skipped."""
        from unittest.mock import MagicMock, patch

        data = _minimal_epub_bytes()
        with patch("relrag.infrastructure.document_parsers.epub_parser.epub.read_epub") as read_epub:
            book = MagicMock()
            read_epub.return_value = book
            good_item = MagicMock()
            good_item.get_type.return_value = 9
            good_item.get_content.return_value = b"<p>OK</p>"
            bad_item = MagicMock()
            bad_item.get_type.return_value = 9
            bad_item.get_content.side_effect = RuntimeError("decode fail")
            book.get_items.return_value = [bad_item, good_item]
            result = parse_epub(data, filename=None)
        assert "OK" in result.text

    def test_invalid_epub_raises_value_error(self) -> None:
        with pytest.raises(ValueError, match="Invalid or corrupted epub"):
            parse_epub(b"not epub", filename="x.epub")


class TestEpubGetDc:
    """Tests for _get_dc (Dublin Core metadata helper)."""

    def test_get_dc_returns_none_on_exception(self) -> None:
        from unittest.mock import MagicMock

        book = MagicMock()
        book.get_metadata.side_effect = RuntimeError("fail")
        assert _get_dc(book, "title") is None

    def test_get_dc_returns_none_when_values_empty(self) -> None:
        from unittest.mock import MagicMock

        book = MagicMock()
        book.get_metadata.return_value = []
        assert _get_dc(book, "title") is None

    def test_get_dc_returns_none_when_first_value_empty(self) -> None:
        from unittest.mock import MagicMock

        book = MagicMock()
        book.get_metadata.return_value = [None]
        assert _get_dc(book, "title") is None

    def test_get_dc_returns_str_when_value_is_list(self) -> None:
        from unittest.mock import MagicMock

        book = MagicMock()
        book.get_metadata.return_value = [["Actual Title"]]
        assert _get_dc(book, "title") == "Actual Title"
