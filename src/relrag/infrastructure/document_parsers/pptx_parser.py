"""Parser for .pptx (PowerPoint)."""

import io
from pathlib import Path

from pptx import Presentation

from relrag.domain.value_objects import PropertyType
from relrag.infrastructure.document_parsers.base import ParseResult
from relrag.infrastructure.document_parsers.metadata_keys import (
    normalize_value_for_storage,
)


def parse_pptx(data: bytes, filename: str | None = None) -> ParseResult:
    """Extract text from slides and core properties from .pptx."""
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as e:
        raise ValueError(f"Invalid or corrupted pptx file: {e}") from e
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    text = "\n\n".join(parts) if parts else " "
    properties: dict[str, tuple[str, PropertyType]] = {}
    cp = prs.core_properties
    if cp.title:
        properties["title"] = (normalize_value_for_storage(cp.title), PropertyType.STRING)
    if cp.author:
        properties["author"] = (normalize_value_for_storage(cp.author), PropertyType.STRING)
    if cp.created:
        properties["created_date"] = (cp.created.isoformat(), PropertyType.DATE)
    if cp.modified:
        properties["modified_date"] = (cp.modified.isoformat(), PropertyType.DATE)
    if cp.last_modified_by:
        properties["author"] = (normalize_value_for_storage(cp.last_modified_by), PropertyType.STRING)
    properties["page_count"] = (str(len(prs.slides)), PropertyType.INT)
    if filename:
        p = Path(filename)
        properties["source_file_name"] = (normalize_value_for_storage(p.name), PropertyType.STRING)
        properties["source_file_type"] = (normalize_value_for_storage("pptx"), PropertyType.STRING)
    return ParseResult(text=text, properties=properties)
